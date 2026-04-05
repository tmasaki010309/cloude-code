# -*- coding: utf-8 -*-
"""
AI副業セミナー PPT生成スクリプト v2
フォーマット：GenSpark風（白背景・カードレイアウト・スプリットセクション）
フォント：Meiryo UI
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ============================================================
# カラーパレット（GenSpark風）
# ============================================================
NAVY      = RGBColor(0x1E, 0x3A, 0x5F)   # ダークネイビー
BLUE      = RGBColor(0x25, 0x63, 0xEB)   # ブライトブルー
LBLUE     = RGBColor(0xEF, 0xF6, 0xFF)   # 薄ブルー背景
BORDER    = RGBColor(0xCB, 0xD5, 0xE1)   # ボーダーグレー
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BG        = RGBColor(0xF8, 0xFA, 0xFC)   # スライド背景（極薄グレー）
CARD_BG   = RGBColor(0xF1, 0xF5, 0xF9)   # カード背景
DTXT      = RGBColor(0x1E, 0x29, 0x3B)   # ダークテキスト
MTXT      = RGBColor(0x64, 0x74, 0x8B)   # ミディアムテキスト
LTXT      = RGBColor(0x94, 0xA3, 0xB8)   # ライトテキスト
RED       = RGBColor(0xDC, 0x26, 0x26)
GREEN     = RGBColor(0x16, 0xA3, 0x4A)
ORANGE    = RGBColor(0xEA, 0x58, 0x0C)
PURPLE    = RGBColor(0x79, 0x16, 0xE8)
GOLD      = RGBColor(0xD9, 0x77, 0x06)

FONT = "Meiryo UI"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ============================================================
# ヘルパー関数
# ============================================================
def bg_slide(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill, line=None, lw=0.75, rounded=False):
    shp = slide.shapes.add_shape(
        5 if rounded else 1,
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    return shp

def add_text(slide, text, l, t, w, h,
             size=14, bold=False, color=DTXT,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_multiline(slide, lines, l, t, w, h,
                  size=13, color=DTXT, indent_color=BLUE):
    """
    lines: list of {"text": str, "level": int, "bold": bool}
    """
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        lv = item.get("level", 0)
        prefix = ("" if lv == 0 else "  " * lv)
        bullet = ("■ " if lv == 0 else "・")
        run = p.add_run()
        run.text = prefix + bullet + item["text"]
        run.font.name = FONT
        run.font.size = Pt(size - lv * 1.5)
        run.font.bold = item.get("bold", lv == 0)
        run.font.color.rgb = indent_color if lv == 0 else color

def content_header(slide, title, session_label="AI副業セミナー"):
    """左縦ライン＋タイトル、右上にセッション名"""
    add_rect(slide, 0, 0, 13.33, 0.9, WHITE)
    add_rect(slide, 0.25, 0.18, 0.07, 0.56, BLUE)          # 縦ライン
    add_text(slide, title, 0.45, 0.2, 9.5, 0.6,
             size=22, bold=True, color=NAVY)
    add_text(slide, session_label, 10.5, 0.28, 2.7, 0.35,
             size=9, color=LTXT, align=PP_ALIGN.RIGHT)
    add_rect(slide, 0, 0.9, 13.33, 0.04, BORDER)            # 下線

def tag_box(slide, text, l, t, color=BLUE, bg=None, text_color=WHITE):
    """ピル型タグ"""
    tw = len(text) * 0.13 + 0.3
    add_rect(slide, l, t, tw, 0.32, bg if bg else color, rounded=True)
    add_text(slide, text, l + 0.08, t + 0.04, tw - 0.16, 0.26,
             size=10, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    return tw

def card(slide, l, t, w, h, title="", title_color=BLUE,
         body_lines=None, header_bg=None, body_bg=WHITE, border=True):
    """カードコンポーネント"""
    lc = BORDER if border else None
    add_rect(slide, l, t, w, h, body_bg, line=lc, lw=0.5)
    if title:
        add_rect(slide, l, t, w, 0.45, header_bg if header_bg else LBLUE, line=lc, lw=0.5)
        add_rect(slide, l, t, 0.05, 0.45, title_color)
        add_text(slide, title, l + 0.15, t + 0.08, w - 0.25, 0.35,
                 size=13, bold=True, color=title_color)
    if body_lines:
        add_multiline(slide, body_lines, l + 0.12, t + 0.55,
                      w - 0.24, h - 0.65, size=12)

def key_message(slide, title, body, l=0.3, t=1.1, w=12.73):
    """KEY MESSAGEボックス"""
    add_rect(slide, l, t, w, 1.05, LBLUE, line=BLUE, lw=1)
    add_rect(slide, l, t, 0.07, 1.05, BLUE)
    add_text(slide, "KEY MESSAGE", l + 0.2, t + 0.08, 2.5, 0.3,
             size=9, bold=True, color=BLUE, italic=True)
    add_text(slide, title, l + 0.2, t + 0.3, w - 0.4, 0.4,
             size=16, bold=True, color=NAVY)
    add_text(slide, body, l + 0.2, t + 0.68, w - 0.4, 0.35,
             size=11, color=MTXT)

def takeaway_bar(slide, items):
    """下部の持ち帰りバー"""
    add_rect(slide, 0, 6.9, 13.33, 0.6, CARD_BG, line=BORDER, lw=0.5)
    add_rect(slide, 0.3, 7.0, 1.8, 0.38, NAVY, rounded=True)
    add_text(slide, "今日の持ち帰り", 0.35, 7.05, 1.7, 0.3,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x = 2.3
    for item in items:
        add_text(slide, "  " + item, x, 7.0, 3.3, 0.38, size=11, color=DTXT)
        x += 3.5

# ============================================================
# SLIDE 1: タイトル
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg_slide(slide, BG)

# ドット背景を模した薄いグリッド（横線）
for i in range(20):
    add_rect(slide, 0, i * 0.38, 13.33, 0.01, RGBColor(0xE2, 0xE8, 0xF0))

# 中央コンテンツエリア
add_rect(slide, 2.5, 1.5, 8.33, 4.8, WHITE, line=BORDER, lw=0.8)

# タグ（角丸ラベル）
add_rect(slide, 4.0, 1.9, 5.33, 0.42, WHITE, line=NAVY, lw=1, rounded=True)
add_text(slide, "AI SIDE BUSINESS SEMINAR", 4.0, 1.95, 5.33, 0.35,
         size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# メインタイトル
add_text(slide, "AI副業セミナー", 2.6, 2.5, 8.1, 1.2,
         size=48, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# 区切り線
add_rect(slide, 5.7, 3.75, 1.93, 0.06, BLUE)

# サブタイトル
add_text(slide, "〜地方からはじめるAI副業の全体像〜",
         2.6, 3.95, 8.1, 0.55, size=18, color=MTXT, align=PP_ALIGN.CENTER)

# ツール名タグ
tools_x = 3.2
for t_name in ["ChatGPT", "Gemini", "Claude", "GenSpark"]:
    add_rect(slide, tools_x, 4.65, 1.4, 0.35, LBLUE, line=BLUE, lw=0.5, rounded=True)
    add_text(slide, t_name, tools_x + 0.05, 4.7, 1.3, 0.28,
             size=10, color=BLUE, align=PP_ALIGN.CENTER)
    tools_x += 1.55

# フッター
add_text(slide, "2026.04  |  Confidential", 0, 7.1, 13.33, 0.35,
         size=10, color=LTXT, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: アジェンダ
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg_slide(slide, WHITE)
content_header(slide, "本日のアジェンダ：全4パートの構成")

# タイムライン縦線
add_rect(slide, 1.8, 1.1, 0.05, 5.6, BORDER)

agenda_items = [
    ("Opening", "オープニング",         "AI副業とは？なぜ今なのかの概要共有",
     ["AIの現在地", "副業の可能性"], MTXT, RGBColor(0xE2, 0xE8, 0xF0)),
    ("Part 1",  "なぜ今AI副業なのか",   "民主化・先行者利益・人材需要の3視点",
     ["AIツール民主化", "先行者利益", "AI人材需要"], BLUE, LBLUE),
    ("Part 2",  "地方×AIの需要",        "地方市場の現状とAIとの相性・キャリアパス",
     ["市場・現状", "AIとの相性", "キャリアパス"], GREEN, RGBColor(0xDC, 0xFC, 0xE7)),
    ("Part 3",  "副業の種類",           "4種類の副業の難易度・収入・獲得方法",
     ["文章生成", "画像生成", "チャットBot", "コンサル"], ORANGE, RGBColor(0xFF, 0xED, 0xCC)),
    ("Part 4",  "仕事の進め方",         "案件獲得から納品までの全ステップ",
     ["案件獲得", "提案", "納品フロー"], PURPLE, RGBColor(0xF3, 0xE8, 0xFF)),
    ("Closing", "まとめ・Q&A",          "今日からできるアクションと質疑応答",
     ["3つのアクション", "Q&A"], MTXT, RGBColor(0xE2, 0xE8, 0xF0)),
]

for i, (label, title, desc, tags, color, tag_bg) in enumerate(agenda_items):
    y = 1.1 + i * 0.94
    # ドット
    add_rect(slide, 1.68, y + 0.08, 0.25, 0.25, color, rounded=True)
    # ラベル
    add_rect(slide, 1.98, y + 0.02, 0.9, 0.3, color if label.startswith("Part") else CARD_BG,
             rounded=True)
    add_text(slide, label, 2.0, y + 0.05, 0.88, 0.26,
             size=9, bold=True, color=WHITE if label.startswith("Part") else MTXT,
             align=PP_ALIGN.CENTER)
    # タイトル
    add_text(slide, title, 3.05, y, 3.5, 0.38,
             size=14, bold=True, color=DTXT)
    # 説明
    add_text(slide, desc, 3.05, y + 0.38, 4.5, 0.45,
             size=11, color=MTXT)
    # タグ
    tx = 7.8
    for tg in tags:
        add_rect(slide, tx, y + 0.1, len(tg) * 0.14 + 0.3, 0.3, tag_bg,
                 line=color, lw=0.5, rounded=True)
        add_text(slide, tg, tx + 0.08, y + 0.13,
                 len(tg) * 0.14 + 0.16, 0.25, size=10, color=color)
        tx += len(tg) * 0.14 + 0.45

# ============================================================
# SLIDE 3: セクション① - スプリットレイアウト
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg_slide(slide, WHITE)
# 左パネル（ネイビー）
add_rect(slide, 0, 0, 4.2, 7.5, NAVY)
# ドット風装飾
for row in range(15):
    for col in range(8):
        add_rect(slide, 0.35 + col * 0.45, 0.35 + row * 0.48, 0.06, 0.06,
                 RGBColor(0x2D, 0x5A, 0x8E))
add_text(slide, "PART", 1.3, 2.6, 1.8, 0.55,
         size=22, bold=True, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)
add_text(slide, "01", 1.0, 3.1, 2.2, 1.1,
         size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(slide, 1.2, 4.3, 1.8, 0.07, RGBColor(0x93, 0xC5, 0xFD))

# 右パネル
add_rect(slide, 4.5, 2.0, 3.5, 0.42, CARD_BG, line=BORDER, lw=0.5, rounded=True)
add_text(slide, "SESSION 1 : WHY AI SIDE BUSINESS", 4.65, 2.08, 3.3, 0.3,
         size=10, color=MTXT, bold=True)
add_text(slide, "なぜ今\nAI副業なのか", 4.5, 2.65, 8.5, 2.0,
         size=52, bold=True, color=NAVY)
add_rect(slide, 4.5, 4.85, 0.08, 1.3, BLUE)
add_text(slide, "AIツールの民主化が\n副業の新しいチャンスを生み出している",
         4.75, 4.9, 8.2, 1.2, size=18, color=MTXT)

# ============================================================
# SLIDE 4: AIツールの民主化（プラスアルファ含む）
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg_slide(slide, WHITE)
content_header(slide, "AIツールの民主化")

key_message(slide,
    "2024年のAI市場規模は約5.5兆円。2030年には30兆円超と予測。",
    "AIツールが急速に普及し、専門知識ゼロでも「使える時代」が到来。今こそ先行する絶好の機会。",
    t=1.1)

# 3つのカード
cards_data = [
    ("AIツールが一般化", BLUE,
     [{"text": "チャット系AIが急速に普及", "level": 0},
      {"text": "ChatGPT・Gemini・Claude", "level": 1},
      {"text": "AIエージェントの登場（自律型）", "level": 0},
      {"text": "GenSpark・Manus：指示するだけで自動実行", "level": 1}]),
    ("初期投資が少ない", GREEN,
     [{"text": "基本はPC・Wi-Fiのみ", "level": 0},
      {"text": "多くのツールは無料〜月数千円", "level": 1},
      {"text": "日本語が使えればOK", "level": 0},
      {"text": "英語スキル不要・翻訳もAIが担う", "level": 1}]),
    ("専門スキルが不要", ORANGE,
     [{"text": "AIへの指示（プロンプト）が主な作業", "level": 0},
      {"text": "試行錯誤しながら誰でも習得可能", "level": 1},
      {"text": "学習コストが劇的に低下", "level": 0},
      {"text": "参入障壁がなくなった副業市場", "level": 1}]),
]
for i, (title, color, lines) in enumerate(cards_data):
    x = 0.3 + i * 4.35
    add_rect(slide, x, 2.35, 4.1, 4.25, WHITE, line=BORDER, lw=0.8)
    add_rect(slide, x, 2.35, 4.1, 0.5, LBLUE, line=BORDER, lw=0.5)
    add_rect(slide, x, 2.35, 0.07, 0.5, color)
    add_text(slide, title, x + 0.18, 2.43, 3.8, 0.38,
             size=14, bold=True, color=color)
    add_multiline(slide, lines, x + 0.15, 2.95, 3.8, 3.5, size=12)

takeaway_bar(slide, ["AIツール3選を今週中に触ってみる", "まず無料プランで始める"])

# ============================================================
# SLIDE 5: 先行者利益 & AI人材の必要性
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg_slide(slide, WHITE)
content_header(slide, "先行者利益 ＆ AI人材の必要性増大")

key_message(slide,
    "AIを仕事に活用できている人はまだ少数派。今が最大のチャンス。",
    "日本のビジネスパーソンのうちAIを業務に活用しているのは約15%（2024年調査）。残り85%が潜在クライアント。",
    t=1.1)

# 左：先行者利益
add_rect(slide, 0.3, 2.35, 6.0, 4.25, WHITE, line=BLUE, lw=1)
add_rect(slide, 0.3, 2.35, 6.0, 0.5, LBLUE, line=BLUE, lw=0.5)
add_rect(slide, 0.3, 2.35, 0.07, 0.5, BLUE)
add_text(slide, "先行者利益", 0.5, 2.43, 5.7, 0.38, size=15, bold=True, color=BLUE)

lines_l = [
    {"text": "AI案件の需要が急拡大中", "level": 0, "bold": True},
    {"text": "中小企業のAI導入支援ニーズが爆発的に増加", "level": 1},
    {"text": "まだ対応できる人材が圧倒的に少ない", "level": 0, "bold": True},
    {"text": "競合が少ない今が参入の絶好タイミング", "level": 1},
    {"text": "副業としてのAI活用はまだ黎明期", "level": 0, "bold": True},
    {"text": "早く始めるほど実績・口コミが積み上がる", "level": 1},
]
add_multiline(slide, lines_l, 0.45, 2.95, 5.7, 3.5, size=12)

# 右：AI人材の必要性
add_rect(slide, 7.0, 2.35, 6.0, 4.25, WHITE, line=ORANGE, lw=1)
add_rect(slide, 7.0, 2.35, 6.0, 0.5, RGBColor(0xFF, 0xF7, 0xED), line=ORANGE, lw=0.5)
add_rect(slide, 7.0, 2.35, 0.07, 0.5, ORANGE)
add_text(slide, "AI人材の必要性増大", 7.2, 2.43, 5.7, 0.38, size=15, bold=True, color=ORANGE)

lines_r = [
    {"text": "AIに代替される仕事が増加", "level": 0, "bold": True},
    {"text": "ルーティン業務・単純作業の自動化が加速", "level": 1},
    {"text": "AIを使いこなす人材の時給が上昇", "level": 0, "bold": True},
    {"text": "同じ時間でより多くの成果を生む人材へ", "level": 1},
    {"text": "企業のAI活用支援ニーズが高まる一方", "level": 0, "bold": True},
    {"text": "外部コンサルタントへの需要が急増", "level": 1},
]
add_multiline(slide, lines_r, 7.15, 2.95, 5.7, 3.5, size=12)

# ============================================================
# SLIDE 6: セクション② 地方×AIの需要
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg_slide(slide, WHITE)
add_rect(slide, 0, 0, 4.2, 7.5, RGBColor(0x06, 0x4E, 0x3B))
for row in range(15):
    for col in range(8):
        add_rect(slide, 0.35 + col * 0.45, 0.35 + row * 0.48, 0.06, 0.06,
                 RGBColor(0x10, 0x7A, 0x56))
add_text(slide, "PART", 1.3, 2.6, 1.8, 0.55, size=22, bold=True,
         color=RGBColor(0x6E, 0xE7, 0xB7), align=PP_ALIGN.CENTER)
add_text(slide, "02", 1.0, 3.1, 2.2, 1.1, size=64, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_rect(slide, 1.2, 4.3, 1.8, 0.07, RGBColor(0x6E, 0xE7, 0xB7))

add_rect(slide, 4.5, 2.0, 3.5, 0.42, CARD_BG, line=BORDER, lw=0.5, rounded=True)
add_text(slide, "SESSION 2 : LOCAL x AI DEMAND", 4.65, 2.08, 3.3, 0.3,
         size=10, color=MTXT, bold=True)
add_text(slide, "地方×AIの\n需要", 4.5, 2.65, 8.5, 2.0,
         size=52, bold=True, color=NAVY)
add_rect(slide, 4.5, 4.85, 0.08, 1.3, GREEN)
add_text(slide, "地方だからこそ競合が少なく\nAI人材の需要が高い「ブルーオーシャン」",
         4.75, 4.9, 8.2, 1.2, size=18, color=MTXT)

# ============================================================
# SLIDE 7: 地方の市場・現状
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg_slide(slide, WHITE)
content_header(slide, "地方の市場・現状")

# 左：数値データ（プラスアルファ）
add_rect(slide, 0.3, 1.1, 5.9, 5.7, WHITE, line=BORDER, lw=0.8)
add_rect(slide, 0.3, 1.1, 5.9, 0.45, LBLUE, line=BORDER, lw=0.5)
add_rect(slide, 0.3, 1.1, 0.07, 0.45, BLUE)
add_text(slide, "地方が抱える課題（データ）", 0.5, 1.17, 5.6, 0.35,
         size=13, bold=True, color=BLUE)

stats = [
    ("2030年までに", "644万人", "の労働力不足が予測（地方中心）"),
    ("地方企業のAI活用率", "わずか8%", "（都市部28%との大きな格差）"),
    ("AI人材の", "9割以上", "が東京・大阪圏に集中"),
]
for i, (pre, num, post) in enumerate(stats):
    y = 1.7 + i * 1.55
    add_rect(slide, 0.45, y, 5.6, 1.35, CARD_BG, line=BORDER, lw=0.5)
    add_text(slide, pre, 0.65, y + 0.1, 5.2, 0.35, size=12, color=MTXT)
    add_text(slide, num, 0.65, y + 0.42, 5.2, 0.55,
             size=28, bold=True, color=NAVY)
    add_text(slide, post, 0.65, y + 0.95, 5.2, 0.35, size=12, color=MTXT)

# 右：深刻な業種
add_rect(slide, 6.5, 1.1, 6.5, 5.7, WHITE, line=BORDER, lw=0.8)
add_rect(slide, 6.5, 1.1, 6.5, 0.45, RGBColor(0xDC, 0xFC, 0xE7), line=BORDER, lw=0.5)
add_rect(slide, 6.5, 1.1, 0.07, 0.45, GREEN)
add_text(slide, "特に深刻な業種と潜在ニーズ", 6.7, 1.17, 6.2, 0.35,
         size=13, bold=True, color=GREEN)

industries = [
    ("運輸・物流", "配車最適化・請求書自動処理・ドライバー不足対応"),
    ("建設・土木", "工程管理・安全報告書作成・図面説明の自動化"),
    ("農林水産", "需要予測・販路開拓・補助金申請書作成支援"),
    ("医療・介護", "記録作成支援・スタッフシフト最適化"),
    ("小売・飲食", "SNS投稿自動化・在庫管理・顧客対応Bot"),
]
for i, (ind, desc) in enumerate(industries):
    y = 1.7 + i * 1.05
    add_rect(slide, 6.65, y, 6.2, 0.95, WHITE, line=BORDER, lw=0.5)
    tag_box(slide, ind, 6.8, y + 0.08, color=GREEN)
    add_text(slide, desc, 6.8, y + 0.5, 6.0, 0.38, size=11, color=MTXT)

takeaway_bar(slide, ["地元の商工会・同業者組合に話を聞いてみる", "AI活用の困りごとをリサーチする"])

# ============================================================
# SLIDE 8: AIキャリアパス
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg_slide(slide, WHITE)
content_header(slide, "地方でのAIキャリアパス")

key_message(slide,
    "まずは「できること」から始め、実績を積み上げながら単価を上げていく。",
    "初期投資はPC・Wi-Fiのみ。完全オンライン対応可能。地域コミュニティ（商工会・同友会）の口コミが最大の武器になる。",
    t=1.1)

steps = [
    ("STEP 1\n初動期", "AI記事作成", BLUE,
     ["クラウドソーシングで実績0から開始",
      "SEO記事・Youtube台本で月1〜3万円",
      "まずはプロフィールと実績づくり"]),
    ("STEP 2\n成長期", "AIチャットBot開発", GREEN,
     ["地元企業への直接営業開始",
      "FAQ Bot・問い合わせ自動化で月5〜20万",
      "商工会・同友会でのデモ展示が有効"]),
    ("STEP 3\n展開期", "AIコンサル", ORANGE,
     ["既存クライアントからの紹介で受注拡大",
      "経営計画書×AI・業務改善で月30万超",
      "オフライン営業がメイン・地域の顔に"]),
]
for i, (step, title, color, points) in enumerate(steps):
    x = 0.3 + i * 4.35
    add_rect(slide, x, 2.35, 4.1, 4.65, WHITE, line=color, lw=1.5)
    add_rect(slide, x, 2.35, 4.1, 0.9, color)
    add_text(slide, step, x + 0.15, 2.42, 3.8, 0.82,
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.15, 3.38, 3.8, 0.55,
             size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
    for j, pt in enumerate(points):
        add_rect(slide, x + 0.2, 4.05 + j * 0.9, 3.7, 0.78,
                 CARD_BG, line=BORDER, lw=0.5)
        add_text(slide, "  ▷  " + pt, x + 0.25, 4.12 + j * 0.9, 3.6, 0.65,
                 size=12, color=DTXT, wrap=True)
    if i < 2:
        add_text(slide, "→", x + 4.0, 4.2, 0.4, 0.5,
                 size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9: セクション③
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg_slide(slide, WHITE)
add_rect(slide, 0, 0, 4.2, 7.5, RGBColor(0x43, 0x18, 0x07))
for row in range(15):
    for col in range(8):
        add_rect(slide, 0.35 + col * 0.45, 0.35 + row * 0.48, 0.06, 0.06,
                 RGBColor(0x7C, 0x2D, 0x12))
add_text(slide, "PART", 1.3, 2.6, 1.8, 0.55, size=22, bold=True,
         color=RGBColor(0xFD, 0xBA, 0x74), align=PP_ALIGN.CENTER)
add_text(slide, "03", 1.0, 3.1, 2.2, 1.1, size=64, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_rect(slide, 1.2, 4.3, 1.8, 0.07, RGBColor(0xFD, 0xBA, 0x74))
add_rect(slide, 4.5, 2.0, 3.0, 0.42, CARD_BG, line=BORDER, lw=0.5, rounded=True)
add_text(slide, "SESSION 3 : TYPE OF AI SIDE BUSINESS", 4.65, 2.08, 2.8, 0.3,
         size=10, color=MTXT, bold=True)
add_text(slide, "副業の種類", 4.5, 2.65, 8.5, 1.2, size=52, bold=True, color=NAVY)
add_rect(slide, 4.5, 4.85, 0.08, 1.3, ORANGE)
add_text(slide, "難易度・収入目安・検索キーワードで\n選ぶ4つのAI副業",
         4.75, 4.9, 8.2, 1.2, size=18, color=MTXT)

# ============================================================
# SLIDE 10: 副業の種類 比較カード
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg_slide(slide, WHITE)
content_header(slide, "副業の種類　比較一覧")

side_jobs = [
    ("AI文章生成", "難易度：低", BLUE, "★☆☆",
     "500〜5,000円/記事",
     "SEO記事作成代行・Youtube台本・メルマガ文章",
     "「記事作成」「ライティング」",
     ["今すぐ始められる入門副業", "クラウドワークスで案件多数"]),
    ("AI画像生成", "難易度：低", GREEN, "★☆☆",
     "1,000〜10,000円/件",
     "SNS素材・EC商品画像・バナー・アイコン作成",
     "「画像作成」「デザイン」",
     ["Midjourney・DALL-E活用", "量産できるため時給が高い"]),
    ("AIチャットBot開発", "難易度：中", ORANGE, "★★☆",
     "5,000〜50,000円/件",
     "FAQ自動化・社内Bot・問い合わせ対応自動化",
     "「チャットボット」「Dify」",
     ["Dify・Copilot Studio活用", "地方企業の需要が特に高い"]),
    ("AIコンサル", "難易度：高", PURPLE, "★★★",
     "月30万円〜（要相談）",
     "AI導入支援・経営計画書作成・業務改善提案",
     "オフライン営業・紹介メイン",
     ["Botからの支援拡大が王道", "高単価・長期契約が期待できる"]),
]
for i, (title, level, color, stars, income, work, search, plus) in enumerate(side_jobs):
    x = 0.25 + i * 3.27
    add_rect(slide, x, 1.1, 3.05, 6.1, WHITE, line=color, lw=1.2)
    add_rect(slide, x, 1.1, 3.05, 0.6, color)
    add_text(slide, title, x + 0.12, 1.17, 2.8, 0.45,
             size=15, bold=True, color=WHITE)
    # 難易度
    add_rect(slide, x + 0.12, 1.82, 1.3, 0.32,
             RGBColor(0xFF, 0xFF, 0xFF), line=color, lw=0.5, rounded=True)
    add_text(slide, level, x + 0.18, 1.86, 1.2, 0.25, size=10, color=color)
    add_text(slide, stars, x + 1.55, 1.86, 1.4, 0.25, size=12, color=GOLD, bold=True)
    # 収入
    add_rect(slide, x + 0.12, 2.28, 2.8, 0.75, CARD_BG, line=BORDER, lw=0.5)
    add_text(slide, "収入目安", x + 0.2, 2.33, 2.6, 0.28, size=9, color=MTXT)
    add_text(slide, income, x + 0.2, 2.57, 2.6, 0.38, size=12, bold=True, color=DTXT)
    # 仕事内容
    add_text(slide, "仕事内容", x + 0.12, 3.15, 2.8, 0.28, size=9, bold=True, color=MTXT)
    add_text(slide, work, x + 0.12, 3.38, 2.8, 0.65, size=11, color=DTXT, wrap=True)
    # 検索方法
    add_rect(slide, x + 0.12, 4.15, 2.8, 0.5, LBLUE, line=BLUE, lw=0.5)
    add_text(slide, "検索ワード：" + search, x + 0.18, 4.22, 2.7, 0.38,
             size=10, color=BLUE, wrap=True)
    # プラスアルファTips
    for j, tip in enumerate(plus):
        add_rect(slide, x + 0.12, 4.8 + j * 0.6, 2.8, 0.5,
                 CARD_BG, line=BORDER, lw=0.5)
        add_text(slide, "  ✓  " + tip, x + 0.15, 4.87 + j * 0.6, 2.75, 0.38,
                 size=10, color=DTXT, wrap=True)

takeaway_bar(slide, ["まずAI文章生成・AI画像生成どちらかに絞る", "今週中にクラウドワークスに登録する"])

# ============================================================
# SLIDE 11: 収益シミュレーション（プラスアルファ）
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg_slide(slide, WHITE)
content_header(slide, "収益シミュレーション　〜現実的な月収モデル〜")

key_message(slide,
    "副業初月から月3万円は現実的。6ヶ月で月10万円超も十分狙える。",
    "実績が積み上がるほど単価が上昇。クラウドソーシング→地域直接営業→紹介の流れが収益を安定させる。",
    t=1.1)

# シミュレーション表
headers = ["フェーズ", "期間目安", "主な副業", "月収目安", "ポイント"]
col_w   = [1.8, 1.5, 2.5, 1.8, 4.4]
col_x   = [0.3]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w + 0.06)

for j, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
    add_rect(slide, x, 2.35, w, 0.48, NAVY)
    add_text(slide, h, x + 0.08, 2.42, w - 0.16, 0.35,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("PHASE 1\n始動期", "1〜2ヶ月",  "AI文章生成・画像生成", "月1〜3万円",
     "クラウドワークス登録→実績0からの受注訓練", BLUE),
    ("PHASE 2\n成長期", "3〜4ヶ月",  "AI文章＋Bot開発",     "月5〜10万円",
     "リピーター獲得・地元企業へ初アプローチ開始",  GREEN),
    ("PHASE 3\n展開期", "5〜6ヶ月",  "Bot開発＋コンサル",   "月10〜30万円",
     "紹介案件が増え始める・単価交渉が可能に",     ORANGE),
    ("PHASE 4\n確立期", "7ヶ月以降", "AIコンサル中心",      "月30万円超",
     "長期契約・顧問契約で安定収入を確立",          PURPLE),
]
for i, (phase, period, job, income, point, color) in enumerate(rows):
    y = 2.9 + i * 1.05
    row_bg = WHITE if i % 2 == 0 else CARD_BG
    row_data = [phase, period, job, income, point]
    for j, (val, x, w) in enumerate(zip(row_data, col_x, col_w)):
        bg_c = row_bg if j > 0 else color
        add_rect(slide, x, y, w, 0.95, bg_c, line=BORDER, lw=0.5)
        add_text(slide, val, x + 0.08, y + 0.1, w - 0.16, 0.78,
                 size=11, bold=(j == 0),
                 color=WHITE if j == 0 else DTXT,
                 align=PP_ALIGN.CENTER if j in [0, 1, 3] else PP_ALIGN.LEFT,
                 wrap=True)

# ============================================================
# SLIDE 12: セクション④
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg_slide(slide, WHITE)
add_rect(slide, 0, 0, 4.2, 7.5, RGBColor(0x2E, 0x10, 0x65))
for row in range(15):
    for col in range(8):
        add_rect(slide, 0.35 + col * 0.45, 0.35 + row * 0.48, 0.06, 0.06,
                 RGBColor(0x5B, 0x21, 0xB6))
add_text(slide, "PART", 1.3, 2.6, 1.8, 0.55, size=22, bold=True,
         color=RGBColor(0xC4, 0xB5, 0xFD), align=PP_ALIGN.CENTER)
add_text(slide, "04", 1.0, 3.1, 2.2, 1.1, size=64, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_rect(slide, 1.2, 4.3, 1.8, 0.07, RGBColor(0xC4, 0xB5, 0xFD))
add_rect(slide, 4.5, 2.0, 3.0, 0.42, CARD_BG, line=BORDER, lw=0.5, rounded=True)
add_text(slide, "SESSION 4 : HOW TO GET PROJECTS", 4.65, 2.08, 2.8, 0.3,
         size=10, color=MTXT, bold=True)
add_text(slide, "仕事の進め方", 4.5, 2.65, 8.5, 1.2, size=48, bold=True, color=NAVY)
add_rect(slide, 4.5, 4.85, 0.08, 1.3, PURPLE)
add_text(slide, "案件獲得から納品まで\n全ステップを解説",
         4.75, 4.9, 8.2, 1.2, size=18, color=MTXT)

# ============================================================
# SLIDE 13: 案件獲得チャネル
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg_slide(slide, WHITE)
content_header(slide, "案件獲得チャネル　〜4つの入口〜")

channels = [
    ("クラウドソーシング", BLUE,
     ["CW（クラウドワークス）", "ココナラ", "ランサーズ"],
     "初心者向け。実績0からでも受注可能。まずここから始める。",
     "初心者向け・今日登録可能"),
    ("紹介・コミュニティ", GREEN,
     ["商工会・同友会", "セミナー参加者", "個人SNS繋がり"],
     "地方では最強の集客手段。信頼ベースで高単価案件に繋がりやすい。",
     "地方最強の集客方法"),
    ("SNS情報発信", ORANGE,
     ["YouTube（実績紹介）", "Instagram・X", "Threads"],
     "継続発信で認知拡大。問い合わせが自動で入る仕組みを作る。",
     "長期的な集客基盤に"),
    ("広告・イベント", PURPLE,
     ["地域セミナー登壇", "Web広告", "商工会主催イベント"],
     "信頼性の高いオフライン展開。地方では特に効果的。",
     "高単価案件獲得の王道"),
]
for i, (title, color, items, desc, tag) in enumerate(channels):
    x = 0.25 + i * 3.27
    add_rect(slide, x, 1.1, 3.05, 5.65, WHITE, line=color, lw=1.2)
    add_rect(slide, x, 1.1, 3.05, 0.55, color)
    add_text(slide, title, x + 0.12, 1.17, 2.8, 0.42, size=14, bold=True, color=WHITE)
    for j, item in enumerate(items):
        add_rect(slide, x + 0.12, 1.78 + j * 0.68, 2.8, 0.58,
                 CARD_BG, line=BORDER, lw=0.5)
        add_text(slide, "  ▷  " + item, x + 0.18, 1.85 + j * 0.68,
                 2.7, 0.45, size=12, color=DTXT)
    add_rect(slide, x + 0.12, 3.88, 2.8, 1.25, LBLUE if color == BLUE else CARD_BG,
             line=BORDER, lw=0.5)
    add_text(slide, desc, x + 0.18, 3.95, 2.72, 1.12,
             size=11, color=MTXT, wrap=True)
    add_rect(slide, x + 0.3, 5.28, 2.4, 0.35, color, rounded=True)
    add_text(slide, tag, x + 0.35, 5.33, 2.3, 0.28,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

takeaway_bar(slide, ["今日：クラウドワークスにプロフィール登録", "今週：地域の商工会・同友会に問い合わせる"])

# ============================================================
# SLIDE 14: 案件フロー
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg_slide(slide, WHITE)
content_header(slide, "案件獲得から納品までの全フロー")

flow = [
    ("01", "案件整理",   BLUE,
     ["オフライン or オンライン", "クラウドソーシングの", "チャット・MTGで詳細確認"]),
    ("02", "提案",       GREEN,
     ["成果物のイメージを", "資料化して提示", "オフライン or MTG"]),
    ("03", "納品物作成", ORANGE,
     ["PCを活用して制作", "AIツールで効率化", "品質チェックを忘れずに"]),
    ("04", "修正",       PURPLE,
     ["フィードバックを", "丁寧に反映", "追加費用の相談も明確に"]),
    ("05", "納品",       RGBColor(0xBE, 0x18, 0x5D),
     ["クラウドソーシング", "or 取り決めた", "格納先へ提出"]),
]
for i, (num, title, color, desc_lines) in enumerate(flow):
    x = 0.25 + i * 2.62
    add_rect(slide, x, 1.1, 2.45, 5.65, WHITE, line=color, lw=1.2)
    add_rect(slide, x, 1.1, 2.45, 0.95, color)
    add_text(slide, num, x + 0.12, 1.15, 0.5, 0.5,
             size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(slide, title, x + 0.12, 1.55, 2.2, 0.45,
             size=16, bold=True, color=WHITE)
    for j, line in enumerate(desc_lines):
        add_rect(slide, x + 0.15, 2.2 + j * 1.1, 2.15, 0.95,
                 CARD_BG, line=BORDER, lw=0.5)
        add_text(slide, line, x + 0.22, 2.3 + j * 1.1, 2.05, 0.75,
                 size=12, color=DTXT, wrap=True, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(slide, "→", x + 2.35, 3.5, 0.35, 0.45,
                 size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ポイント補足
add_rect(slide, 0.3, 6.9, 12.73, 0.42, LBLUE, line=BLUE, lw=0.5)
add_text(slide, "  ■  完全オンライン対応可能（請求書・契約書もDX化）　■  初期費用：PC・Wi-Fiのみ",
         0.4, 6.95, 12.5, 0.33, size=11, color=BLUE, bold=True)

# ============================================================
# SLIDE 15: よくある失敗と対策（プラスアルファ）
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg_slide(slide, WHITE)
content_header(slide, "よくある失敗パターンと対策")

key_message(slide,
    "多くの人が最初の3ヶ月で諦める。原因は「完璧主義」と「集客不足」。",
    "失敗パターンを事前に知っておくだけで、継続率は大きく変わる。",
    t=1.1)

failures = [
    ("完璧になってから始める",
     "ツールを覚えてから・スキルが上がってから…と先延ばしにして参入時期を逃す",
     "まず受注してみる。品質は後から上げられる。80点で納品→改善サイクルを回す",
     RED),
    ("クラウドソーシングだけに依存",
     "価格競争が激しく、単価が上がりにくい。時間当たり収益が伸びない",
     "3ヶ月を目処にSNS発信と地域営業を並行する。紹介案件を早期に作る",
     ORANGE),
    ("1人でやろうとしすぎる",
     "わからないことで詰まり、時間を浪費。孤独感からモチベーションが低下",
     "AIコミュニティ・副業仲間を作る。同じ境遇の人と情報交換する習慣をつける",
     PURPLE),
]
for i, (title, problem, solution, color) in enumerate(failures):
    x = 0.3 + i * 4.35
    add_rect(slide, x, 2.35, 4.1, 4.3, WHITE, line=color, lw=1)
    add_rect(slide, x, 2.35, 4.1, 0.5, color)
    add_text(slide, "NG " + title, x + 0.12, 2.43, 3.85, 0.38,
             size=13, bold=True, color=WHITE)
    add_rect(slide, x + 0.12, 3.0, 3.85, 1.3, RGBColor(0xFF, 0xF1, 0xF0), line=RED, lw=0.5)
    add_text(slide, "問題", x + 0.18, 3.05, 0.6, 0.28, size=9, bold=True, color=RED)
    add_text(slide, problem, x + 0.18, 3.3, 3.72, 0.92, size=11, color=DTXT, wrap=True)
    add_rect(slide, x + 0.12, 4.45, 3.85, 1.9, RGBColor(0xF0, 0xFF, 0xF4), line=GREEN, lw=0.5)
    add_text(slide, "対策", x + 0.18, 4.5, 0.6, 0.28, size=9, bold=True, color=GREEN)
    add_text(slide, solution, x + 0.18, 4.75, 3.72, 1.45, size=11, color=DTXT, wrap=True)

# ============================================================
# SLIDE 16: まとめ・3つのアクション
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg_slide(slide, WHITE)
content_header(slide, "まとめ　〜今日からできる3つのアクション〜")

# サマリーカード 4つ
summaries = [
    ("AIツールの民主化", "専門スキル不要・低投資で今すぐ始められる時代", BLUE),
    ("地方×AIの優位性", "競合少・需要高・コミュニティが最大の武器", GREEN),
    ("副業の種類", "文章→画像→Bot→コンサルとステップアップ", ORANGE),
    ("仕事の進め方", "クラウドソーシングで実績を積み、地域直営業へ", PURPLE),
]
for i, (title, desc, color) in enumerate(summaries):
    x = 0.25 + i * 3.27
    add_rect(slide, x, 1.1, 3.05, 1.6, WHITE, line=color, lw=1)
    add_rect(slide, x, 1.1, 0.07, 1.6, color)
    add_text(slide, title, x + 0.2, 1.18, 2.75, 0.42, size=13, bold=True, color=color)
    add_text(slide, desc, x + 0.2, 1.6, 2.75, 0.95, size=11, color=MTXT, wrap=True)

# 3つのアクション
add_rect(slide, 0.3, 3.0, 12.73, 0.42, NAVY)
add_text(slide, "今日からできる　3つのアクション", 0.5, 3.07, 12.3, 0.32,
         size=14, bold=True, color=WHITE)

actions = [
    ("ACTION 1\n今日", "クラウドワークスに登録し、プロフィールを完成させる",
     "https://crowdworks.jp　→「AIライティング」で案件検索", BLUE),
    ("ACTION 2\n今週", "ChatGPT・Claude・Geminiのいずれかを1週間使い込む",
     "まず1つに絞って毎日触れる。プロンプトの感覚を掴む", GREEN),
    ("ACTION 3\n今月", "地元の商工会・同友会に連絡し、AI活用の困りごとを聞く",
     "名刺代わりの「AIデモ」を1つ準備しておくと効果的", ORANGE),
]
for i, (num, title, desc, color) in enumerate(actions):
    x = 0.3 + i * 4.35
    add_rect(slide, x, 3.6, 4.1, 3.0, WHITE, line=color, lw=1.2)
    add_rect(slide, x, 3.6, 4.1, 0.65, color)
    add_text(slide, num, x + 0.15, 3.65, 1.0, 0.58,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.15, 4.4, 3.8, 0.65,
             size=13, bold=True, color=color, wrap=True)
    add_rect(slide, x + 0.15, 5.2, 3.8, 1.2, CARD_BG, line=BORDER, lw=0.5)
    add_text(slide, desc, x + 0.22, 5.28, 3.65, 1.05,
             size=11, color=MTXT, wrap=True)

add_text(slide, "まずは「完璧」を目指さず、小さく始めてみることが最大のコツです。",
         0.3, 6.7, 12.73, 0.5, size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ============================================================
# 保存
# ============================================================
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "下書き_v2.pptx")
prs.save(out_path)
print("DONE: " + out_path)
print("Slides: " + str(len(prs.slides)))
